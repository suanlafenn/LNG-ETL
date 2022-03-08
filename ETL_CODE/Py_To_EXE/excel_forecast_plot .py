#!/usr/bin/env python
# coding: utf-8

# In[11]:


import numpy as np
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches 
import io


# In[12]:


prs = Presentation()
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "Hello, World!"
subtitle.text = "python-pptx was here!"

prs.save('test.pptx')


# In[13]:


np.random.seed(5)
x = np.arange(1, 101)
y = 20 + 3 * x + np.random.normal(0, 60, 100)
plt.plot(x, y, "o")
plt.show()
image_stream = io.BytesIO()
plt.savefig(image_stream)


prs = Presentation()
blank_slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(blank_slide_layout)
x = Inches(1)
y = Inches(5)
z = Inches(5.5)
pic = slide.shapes.add_picture(image_stream, x, y, z)
prs.save('test.pptx')

